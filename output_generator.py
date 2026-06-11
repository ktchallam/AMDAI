import os
import re
from datetime import datetime
from typing import List, Dict, Any
from pathlib import Path
from loguru import logger

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from config import get_config
from answer_agent import QAPair
from validation_agent import ValidationResult

class CheatSheetGenerator:
    def __init__(self, output_dir: str = "./output"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        company: str,
        quarter: str,
        validated_pairs: List[ValidationResult],
        data_gaps: List[Any],
        kpi_delta: Dict[str, Any],
        overall_score: float = 0.0
    ) -> str:
        """Generates a markdown CEO/CFO preparation cheat sheet brief."""
        logger.info(f"Generating Brief Cheat Sheet for {company} ({quarter})...")
        
        # 1. Build KPI table
        kpis = kpi_delta.get("kpis", [])
        kpi_table_rows = []
        kpi_table_rows.append("| Metric | Guided | Actual | Delta % | Status |")
        kpi_table_rows.append("| --- | --- | --- | --- | --- |")
        
        for k in kpis:
            status = "MISS" if k["is_miss"] else "BEAT"
            # If actual matches guided, status is "IN LINE"
            if abs(k["delta_pct"]) < 0.5:
                status = "IN LINE"
            kpi_table_rows.append(
                f"| {k['metric']} | {k['guided_value']}{k['unit']} | "
                f"{k['actual_value']}{k['unit']} | {k['delta_pct']:.2f}% | {status} |"
            )
            
        kpi_table_str = "\n".join(kpi_table_rows)
        
        # 2. Build Predicted Questions
        q_sections = []
        for idx, res in enumerate(validated_pairs):
            qa = res.qa_pair
            diff_pct = int(qa.adversarial_score * 100)
            sources = ", ".join(qa.source_quarters) if qa.source_quarters else "None"
            q_sections.append(
                f"### Q{idx+1}: {qa.question}\n\n"
                f"**Topic**: {qa.topic} | **Difficulty**: {diff_pct}% | **Historical Quarters**: {sources}\n\n"
                f"> **Suggested Answer**:\n"
                f"> {res.final_answer}\n\n"
                f"**Why Tough**: {qa.why_tough}\n"
            )
            
        questions_str = "\n---\n".join(q_sections)
        
        # 3. Build Data Gaps
        gap_sections = []
        if not data_gaps:
            gap_sections.append("*No data gaps identified. All anticipated questions are covered by the current disclosures.*")
        else:
            for idx, gap in enumerate(data_gaps):
                gap_sections.append(
                    f"### Gap {idx+1}: {gap.question}\n"
                    f"**Topic**: {gap.topic} | **Adversarial Score**: {int(gap.adversarial_score*100)}%\n"
                    f"**Reason**: {gap.gap_reason}\n"
                    f"*Action Required*: **CFO must prepare response details manually.**\n"
                )
        gaps_str = "\n".join(gap_sections)
        
        # Markdown layout
        md_content = f"""# {company} {quarter} Earnings Call — CEO/CFO Preparation Brief

## Quality Metrics
- **Overall Q&A Quality Score**: {overall_score:.2f}/1.0
- **Anticipated Answerable Questions**: {len(validated_pairs)}
- **Contextual Data Gaps**: {len(data_gaps)}

## KPI Summary vs Guidance
{kpi_table_str}

## Predicted Questions — Ranked by Difficulty
{questions_str}

## Data Gaps — Manual Preparation Required
{gaps_str}
"""
        
        filename = f"{company.replace(' ', '_')}_{quarter}_brief.md"
        filepath = self.output_dir / filename
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(md_content)
            
        logger.info(f"Brief Cheat Sheet saved to {filepath}")
        return str(filepath)

class PresentationGenerator:
    def __init__(self, output_dir: str = "./output"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Styling tokens
        self.bg_color = RGBColor(0x1a, 0x1f, 0x36)  # Dark Navy
        self.text_white = RGBColor(0xff, 0xff, 0xff)
        self.accent_blue = RGBColor(0x4f, 0x9c, 0xf9) # Accent Blue
        self.text_gray = RGBColor(0xab, 0xaf, 0xb8)
        
        # KPI Colors
        self.color_beat = RGBColor(0x2e, 0x7d, 0x32) # Premium Dark Green
        self.color_miss = RGBColor(0xc6, 0x28, 0x28) # Premium Dark Red
        self.color_inline = RGBColor(0x42, 0x42, 0x42) # Muted Gray

    def _apply_dark_background(self, slide):
        """Fills slide background with premium dark navy."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color

    def _add_slide_title(self, slide, text: str):
        """Helper to create a unified header across content slides."""
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.text_white
        p.alignment = PP_ALIGN.LEFT

    def generate(
        self,
        company: str,
        quarter: str,
        validated_pairs: List[ValidationResult],
        kpi_delta: Dict[str, Any],
        topics: List[Dict[str, Any]],
        data_gaps: List[Any]
    ) -> str:
        """Generates a professional 16:9 presentation deck (.pptx)."""
        logger.info(f"Generating Slides Presentation for {company} ({quarter})...")
        
        prs = Presentation()
        # Set 16:9 widescreen dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Use a blank layout (usually layout index 6 is blank in default template)
        blank_layout = prs.slide_layouts[6]
        
        # ==========================================
        # SLIDE 1: TITLE SLIDE
        # ==========================================
        slide1 = prs.slides.add_slide(blank_layout)
        self._apply_dark_background(slide1)
        
        # Title text block
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = f"{company} — {quarter}"
        p.font.name = "Calibri"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.text_white
        p.alignment = PP_ALIGN.LEFT
        
        p2 = tf.add_paragraph()
        p2.text = "Earnings Call Q&A & Performance Analysis"
        p2.font.name = "Calibri"
        p2.font.size = Pt(28)
        p2.font.color.rgb = self.accent_blue
        p2.font.bold = True
        p2.alignment = PP_ALIGN.LEFT
        
        p3 = tf.add_paragraph()
        p3.text = f"Investor Relations Team Briefing | {datetime.now().strftime('%B %d, %Y')}"
        p3.font.name = "Calibri"
        p3.font.size = Pt(16)
        p3.font.color.rgb = self.text_gray
        p3.alignment = PP_ALIGN.LEFT

        # ==========================================
        # SLIDE 2: KPI DASHBOARD
        # ==========================================
        slide2 = prs.slides.add_slide(blank_layout)
        self._apply_dark_background(slide2)
        self._add_slide_title(slide2, "KPI Performance vs Guidance")
        
        # Add a table
        kpis = kpi_delta.get("kpis", [])
        rows = len(kpis) + 1
        cols = 5
        
        left = Inches(0.75)
        top = Inches(1.8)
        width = Inches(11.83)
        height = Inches(4.5)
        
        table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Columns widths
        table.columns[0].width = Inches(3.5)  # Metric
        table.columns[1].width = Inches(2.0)  # Guided
        table.columns[2].width = Inches(2.0)  # Actual
        table.columns[3].width = Inches(2.0)  # Delta
        table.columns[4].width = Inches(2.33) # Status
        
        headers = ["Financial Metric", "Guidance", "Actual Reported", "Variance %", "Status"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.accent_blue
            
            # Format header text
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.text_white
            p.alignment = PP_ALIGN.CENTER
            
        for row_idx, k in enumerate(kpis):
            status = "MISS" if k["is_miss"] else "BEAT"
            if abs(k["delta_pct"]) < 0.5:
                status = "IN LINE"
                
            row_data = [
                k["metric"],
                f"{k['guided_value']}{k['unit']}",
                f"{k['actual_value']}{k['unit']}",
                f"{k['delta_pct']:.2f}%",
                status
            ]
            
            # Select background status color
            bg_color = self.color_beat if status == "BEAT" else (self.color_miss if status == "MISS" else self.color_inline)
            
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                
                # Format cell text
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Calibri"
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_white
                if col_idx == 0:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.CENTER

        # ==========================================
        # SLIDE 3: ANALYST SENTIMENT
        # ==========================================
        slide3 = prs.slides.add_slide(blank_layout)
        self._apply_dark_background(slide3)
        self._add_slide_title(slide3, "Historical Analyst Sentiment Trends")
        
        # Display top topics
        rows = min(8, len(topics)) + 1
        cols = 4
        
        table_shape3 = slide3.shapes.add_table(rows, cols, left, top, width, height)
        table3 = table_shape3.table
        table3.columns[0].width = Inches(4.5)  # Topic
        table3.columns[1].width = Inches(2.2)  # Recurrence
        table3.columns[2].width = Inches(2.2)  # Avg Sentiment
        table3.columns[3].width = Inches(2.93) # Trend
        
        headers3 = ["Mined Concern Topic", "Quarters Recurrence", "Avg Sentiment", "Sentiment Outlook"]
        for col_idx, text in enumerate(headers3):
            cell = table3.cell(0, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.accent_blue
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.text_white
            p.alignment = PP_ALIGN.CENTER
            
        for row_idx, t in enumerate(topics[:8]):
            recurrence_pct = f"{int(t['recurrence'] * 100)}%"
            sentiment_score = f"{t['avg_sentiment']:.2f}"
            
            trend = "Bullish/Positive" if t["avg_sentiment"] > 0.1 else ("Bearish/Negative" if t["avg_sentiment"] < -0.1 else "Neutral")
            trend_color = self.color_beat if t["avg_sentiment"] > 0.1 else (self.color_miss if t["avg_sentiment"] < -0.1 else self.color_inline)
            
            row_data = [t["label"], recurrence_pct, sentiment_score, trend]
            
            for col_idx, val in enumerate(row_data):
                cell = table3.cell(row_idx + 1, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = trend_color if col_idx == 3 else self.bg_color
                
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Calibri"
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_white
                if col_idx == 0:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.CENTER

        # ==========================================
        # SLIDES 4-N: TOP 5 QUESTIONS (ONE PER SLIDE)
        # ==========================================
        for idx, res in enumerate(validated_pairs[:5]):
            qa = res.qa_pair
            slide_q = prs.slides.add_slide(blank_layout)
            self._apply_dark_background(slide_q)
            self._add_slide_title(slide_q, f"Anticipated: {qa.topic}")
            
            # Add Question Card text box
            q_box = slide_q.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(1.2))
            tf_q = q_box.text_frame
            tf_q.word_wrap = True
            tf_q.margin_left = tf_q.margin_top = tf_q.margin_bottom = tf_q.margin_right = 0
            
            p_lbl = tf_q.paragraphs[0]
            p_lbl.text = f"ANALYST QUESTION {idx+1}:"
            p_lbl.font.name = "Calibri"
            p_lbl.font.size = Pt(14)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = self.accent_blue
            
            p_q = tf_q.add_paragraph()
            p_q.text = f"\"{qa.question}\""
            p_q.font.name = "Calibri"
            p_q.font.size = Pt(18)
            p_q.font.italic = True
            p_q.font.bold = True
            p_q.font.color.rgb = self.text_white
            
            # Suggested answer bullets
            ans_box = slide_q.shapes.add_textbox(Inches(0.75), Inches(2.9), Inches(11.83), Inches(3.6))
            tf_a = ans_box.text_frame
            tf_a.word_wrap = True
            tf_a.margin_left = tf_a.margin_top = tf_a.margin_bottom = tf_a.margin_right = 0
            
            p_albl = tf_a.paragraphs[0]
            p_albl.text = "SUGGESTED TALKING POINTS:"
            p_albl.font.name = "Calibri"
            p_albl.font.size = Pt(14)
            p_albl.font.bold = True
            p_albl.font.color.rgb = self.accent_blue
            
            # Split answer to 3 bullet points
            sentences = [s.strip() for s in re.split(r'(?<=[.?!])\s+', res.final_answer) if s.strip()]
            if not sentences:
                sentences = ["No answer talking points generated."]
                
            # Group into max 3 points
            bullets = sentences[:3]
            if len(sentences) > 3:
                # Append remainder to 3rd bullet
                bullets[2] = " ".join(sentences[2:])
                
            for b in bullets:
                p_b = tf_a.add_paragraph()
                p_b.text = f"• {b}"
                p_b.font.name = "Calibri"
                p_b.font.size = Pt(16)
                p_b.font.color.rgb = self.text_white
                p_b.space_after = Pt(10)
                
            # Slide Footer Info
            footer_box = slide_q.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(11.83), Inches(0.5))
            tf_f = footer_box.text_frame
            tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
            p_f = tf_f.paragraphs[0]
            diff_pct = int(qa.adversarial_score * 100)
            quarters_str = ", ".join(qa.source_quarters) if qa.source_quarters else "None"
            p_f.text = f"Difficulty Score: {diff_pct}%  |  Historical Recurrence Quarters: {quarters_str}"
            p_f.font.name = "Calibri"
            p_f.font.size = Pt(12)
            p_f.font.color.rgb = self.text_gray

        # ==========================================
        # LAST SLIDE: APPENDIX (DATA GAPS)
        # ==========================================
        slide_app = prs.slides.add_slide(blank_layout)
        self._apply_dark_background(slide_app)
        self._add_slide_title(slide_app, "Appendix: Identified Data Gaps")
        
        gap_box = slide_app.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8))
        tf_gap = gap_box.text_frame
        tf_gap.word_wrap = True
        tf_gap.margin_left = tf_gap.margin_right = tf_gap.margin_top = tf_gap.margin_bottom = 0
        
        p_desc = tf_gap.paragraphs[0]
        p_desc.text = "The following questions are highly probable but lack sufficient disclosures in the current quarter's filing. IR/CFO must prepare manual answers:"
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(16)
        p_desc.font.color.rgb = self.text_gray
        p_desc.space_after = Pt(15)
        
        if not data_gaps:
            p_none = tf_gap.add_paragraph()
            p_none.text = "• None. All topics are fully covered by the current disclosures."
            p_none.font.name = "Calibri"
            p_none.font.size = Pt(16)
            p_none.font.color.rgb = self.text_white
        else:
            for gap in data_gaps[:5]:  # limit to top 5
                p_g = tf_gap.add_paragraph()
                p_g.text = f"• Q: {gap.question}\n   Reason: {gap.gap_reason}"
                p_g.font.name = "Calibri"
                p_g.font.size = Pt(14)
                p_g.font.color.rgb = self.text_white
                p_g.space_after = Pt(10)
                
        filename_pptx = f"{company.replace(' ', '_')}_{quarter}_deck.pptx"
        filepath_pptx = self.output_dir / filename_pptx
        prs.save(str(filepath_pptx))
        
        logger.info(f"Presentation saved to {filepath_pptx}")
        return str(filepath_pptx)

class OutputOrchestrator:
    def __init__(self, output_dir: str = "./output"):
        self.output_dir = output_dir
        self.cheat_sheet_gen = CheatSheetGenerator(output_dir)
        self.presentation_gen = PresentationGenerator(output_dir)

    def run(
        self,
        company: str,
        quarter: str,
        validation_output: Dict[str, Any],
        topics: List[Dict[str, Any]],
        kpi_delta: Dict[str, Any]
    ) -> Dict[str, str]:
        """Orchestrates cheat sheet and slide deck compilation."""
        validated_pairs = validation_output.get("validated_pairs", [])
        data_gaps = validation_output.get("data_gaps", [])
        overall_score = validation_output.get("overall_quality_score", 0.0)
        
        cheat_sheet_path = self.cheat_sheet_gen.generate(
            company=company,
            quarter=quarter,
            validated_pairs=validated_pairs,
            data_gaps=data_gaps,
            kpi_delta=kpi_delta,
            overall_score=overall_score
        )
        
        pptx_path = self.presentation_gen.generate(
            company=company,
            quarter=quarter,
            validated_pairs=validated_pairs,
            kpi_delta=kpi_delta,
            topics=topics,
            data_gaps=data_gaps
        )
        
        summary = (
            f"Successfully generated preparation package for {company} {quarter}.\n"
            f"1. Executive Briefing Sheet: {cheat_sheet_path}\n"
            f"2. Presentation Slide Deck: {pptx_path}"
        )
        
        return {
            "cheat_sheet_path": cheat_sheet_path,
            "pptx_path": pptx_path,
            "summary": summary
        }
